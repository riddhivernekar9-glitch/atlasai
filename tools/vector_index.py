"""
tools/vector_index.py
─────────────────────
Ingestion and semantic search for AtlasAI.

Storage  : ChromaDB persistent store (replaces flat JSON index)
Embedding: sentence-transformers all-MiniLM-L6-v2  (384-dim, cosine)
Search   : ChromaDB cosine similarity query
"""
from __future__ import annotations

import os
import platform
import threading
from pathlib import Path
from typing import Any, Dict, List

import hashlib

__all__ = [
    "ingest_folder",
    "search_index",
    "_reset_collection",   # exported for test isolation
    "_app_data_dir",
    "_embed",
]

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

ALLOWED_EXTENSIONS = {
    ".pdf", ".txt", ".md",
    ".docx", ".xlsx", ".pptx",
    ".html", ".htm",
}

SKIP_DIR_NAMES = {
    "node_modules", "venv", ".venv", "__pycache__",
    ".git", "dist", "build",
    ".pytest_cache", ".mypy_cache", ".ruff_cache",
}

CHUNK_SIZE = 1000
CHUNK_OVERLAP = 150
CHROMA_COLLECTION = "atlasai_docs"
EMBED_MODEL = "all-MiniLM-L6-v2"


# ---------------------------------------------------------------------------
# OS data directory
# ---------------------------------------------------------------------------

def _app_data_dir() -> Path:
    """Return the OS-appropriate user data directory for AtlasAI."""
    system = platform.system()
    if system == "Darwin":
        return Path.home() / "Library" / "Application Support" / "AtlasAI"
    elif system == "Windows":
        appdata = os.environ.get("APPDATA") or ""
        base = Path(appdata) if appdata else Path.home() / "AppData" / "Roaming"
        return base / "AtlasAI"
    else:
        xdg = os.environ.get("XDG_DATA_HOME") or ""
        base = Path(xdg) if xdg else Path.home() / ".local" / "share"
        return base / "AtlasAI"


# ---------------------------------------------------------------------------
# Singleton: embedding model
# ---------------------------------------------------------------------------

_model: Any = None
_model_lock = threading.Lock()


def _get_model() -> Any:
    global _model
    if _model is None:
        with _model_lock:
            if _model is None:
                from sentence_transformers import SentenceTransformer  # type: ignore
                _model = SentenceTransformer(EMBED_MODEL)
    return _model


def _embed(texts: List[str]) -> List[List[float]]:
    """Return L2-normalised embeddings for a list of strings."""
    model = _get_model()
    vecs = model.encode(
        texts,
        normalize_embeddings=True,
        show_progress_bar=False,
        batch_size=32,
    )
    return vecs.tolist()


# ---------------------------------------------------------------------------
# Singleton: ChromaDB collection
# ---------------------------------------------------------------------------

_collection: Any = None
_collection_lock = threading.Lock()


def _get_collection() -> Any:
    global _collection
    if _collection is None:
        with _collection_lock:
            if _collection is None:
                import chromadb  # type: ignore
                db_path = str(_app_data_dir() / "chroma")
                Path(db_path).mkdir(parents=True, exist_ok=True)
                client = chromadb.PersistentClient(path=db_path)
                _collection = client.get_or_create_collection(
                    name=CHROMA_COLLECTION,
                    metadata={"hnsw:space": "cosine"},
                )
    return _collection


def _reset_collection() -> None:
    """Force the collection singleton to be recreated on next access.
    Called in tests to isolate each test's ChromaDB store."""
    global _collection
    _collection = None


# ---------------------------------------------------------------------------
# Path / file utilities
# ---------------------------------------------------------------------------

def normalize_input_path(raw_path: str) -> Path:
    if not raw_path or not raw_path.strip():
        raise ValueError("No path provided")
    cleaned = raw_path.strip().strip('"').replace("\r", "").replace("\n", "")
    path = Path(cleaned).expanduser()
    try:
        path = path.resolve(strict=True)
    except FileNotFoundError:
        raise ValueError(f"Path does not exist or is not accessible: {cleaned}")
    return path


def _should_index_file(p: Path) -> bool:
    if not p.exists() or not p.is_file():
        return False
    if p.suffix.lower() not in ALLOWED_EXTENSIONS:
        return False
    parts = set(p.parts)
    return not any(name in parts for name in SKIP_DIR_NAMES)


def _iter_indexable_files(path_or_folder: str) -> List[Path]:
    try:
        base = normalize_input_path(path_or_folder)
    except Exception:
        return []
    if base.exists() and base.is_file():
        return [base] if _should_index_file(base) else []
    if not base.exists() or not base.is_dir():
        return []
    return sorted(p for p in base.rglob("*") if _should_index_file(p))


# ---------------------------------------------------------------------------
# File readers
# ---------------------------------------------------------------------------

def _read_text_file(p: Path) -> str:
    try:
        return p.read_text(errors="ignore")
    except Exception:
        return ""


def _read_pdf_file(p: Path) -> str:
    try:
        from server.utils.pdf_reader import read_pdf_text
        return read_pdf_text(p)
    except Exception:
        return ""


def _read_docx_file(p: Path) -> str:
    try:
        from docx import Document  # type: ignore
        doc = Document(str(p))
        return "\n".join(para.text for para in doc.paragraphs if para.text.strip())
    except Exception:
        return ""


def _read_xlsx_file(p: Path) -> str:
    try:
        from openpyxl import load_workbook  # type: ignore
        wb = load_workbook(str(p), read_only=True, data_only=True)
        lines: List[str] = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None and str(c).strip()]
                if cells:
                    lines.append("\t".join(cells))
        return "\n".join(lines)
    except Exception:
        return ""


def _read_pptx_file(p: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore
        prs = Presentation(str(p))
        lines: List[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            lines.append(text)
        return "\n".join(lines)
    except Exception:
        return ""


def _read_html_file(p: Path) -> str:
    try:
        from bs4 import BeautifulSoup  # type: ignore
        html = p.read_text(errors="ignore")
        soup = BeautifulSoup(html, "lxml")
        for tag in soup(["script", "style", "nav", "footer"]):
            tag.decompose()
        return soup.get_text(separator="\n", strip=True)
    except Exception:
        return ""


def _read_file(p: Path) -> str:
    ext = p.suffix.lower()
    if ext == ".pdf":
        return _read_pdf_file(p)
    if ext == ".docx":
        return _read_docx_file(p)
    if ext == ".xlsx":
        return _read_xlsx_file(p)
    if ext == ".pptx":
        return _read_pptx_file(p)
    if ext in (".html", ".htm"):
        return _read_html_file(p)
    return _read_text_file(p)


# ---------------------------------------------------------------------------
# Chunking / hashing
# ---------------------------------------------------------------------------

def _chunk_text(text: str) -> List[str]:
    if not text:
        return []
    chunks: List[str] = []
    start = 0
    n = len(text)
    while start < n:
        end = min(start + CHUNK_SIZE, n)
        chunk = text[start:end]
        if chunk:
            chunks.append(chunk)
        start += max(1, CHUNK_SIZE - CHUNK_OVERLAP)
    return chunks


def _hash_text(text: str) -> str:
    return hashlib.sha256(text.encode("utf-8", errors="ignore")).hexdigest()


# ---------------------------------------------------------------------------
# Document classifier
# ---------------------------------------------------------------------------

def classify_document(path: str, text: str) -> str:
    combined = f"{path}\n{text}".lower()
    if any(x in combined for x in [
        "sponsor", "sponsorship", "employment visa", "residence permit", "work permit"
    ]):
        return "Visa"
    if any(x in combined for x in [
        "salary", "payroll", "wage", "compensation", "allowance"
    ]):
        return "Payroll"
    if any(x in combined for x in [
        "contract", "agreement", "terms and conditions", "probation", "termination"
    ]):
        return "Contract"
    if any(x in combined for x in [
        "insurance", "medical insurance", "coverage", "insured"
    ]):
        return "Insurance"
    if any(x in combined for x in [
        "law", "regulation", "compliance", "policy", "obligation"
    ]):
        return "Compliance"
    return "General"


# ---------------------------------------------------------------------------
# Public API: ingest
# ---------------------------------------------------------------------------

def ingest_folder(folder_path: str) -> Dict[str, Any]:
    """
    Read all supported files under *folder_path*, chunk them, embed with
    all-MiniLM-L6-v2, and upsert into the ChromaDB collection.

    Returns the same dict shape the FastAPI routes expect:
    {"ok": bool, "count": int, "chunks": int, "files": [...]}
    """
    files = _iter_indexable_files(folder_path)
    if not files:
        return {"ok": False, "error": "No valid files found."}

    collection = _get_collection()
    total_chunks = 0
    files_seen = 0

    for p in files:
        content = _read_file(p)
        if not content.strip():
            content = f"Document file: {p.name}"

        doc_type = classify_document(str(p), content)
        chunks = _chunk_text(content) or [f"Document file: {p.name}"]

        chunk_ids = [_hash_text(str(p) + "::" + chunk) for chunk in chunks]
        embeddings = _embed(chunks)
        metadatas = [
            {"relative_path": p.name, "path": str(p), "doc_type": doc_type}
            for _ in chunks
        ]

        # upsert: same ID → overwrite (idempotent re-index)
        collection.upsert(
            ids=chunk_ids,
            embeddings=embeddings,
            documents=chunks,
            metadatas=metadatas,
        )

        total_chunks += len(chunks)
        files_seen += 1

    return {
        "ok": True,
        "folder": str(Path(folder_path).expanduser()),
        "count": files_seen,
        "chunks": total_chunks,
        "files": [str(f) for f in files],
    }


# ---------------------------------------------------------------------------
# Public API: search
# ---------------------------------------------------------------------------

def search_index(query: str, limit: int = 3) -> List[Dict[str, Any]]:
    """
    Embed *query* and return the top *limit* most similar chunks from ChromaDB.

    Each result dict contains:
        relative_path, path, context, score (cosine similarity 0-1), doc_type
    """
    if not query or not query.strip():
        return []

    collection = _get_collection()
    total = collection.count()
    if total == 0:
        return []

    n = min(max(1, int(limit)), total)
    query_embedding = _embed([query.strip()])[0]

    raw = collection.query(
        query_embeddings=[query_embedding],
        n_results=n,
        include=["documents", "metadatas", "distances"],
    )

    ids       = raw["ids"][0]
    documents = raw["documents"][0]
    metadatas = raw["metadatas"][0]
    distances = raw["distances"][0]

    results: List[Dict[str, Any]] = []
    for _, doc, meta, dist in zip(ids, documents, metadatas, distances):
        # ChromaDB cosine distance = 1 − cosine_similarity  →  similarity = 1 − dist
        score = round(float(1.0 - dist), 4)
        results.append({
            "relative_path": meta.get("relative_path", ""),
            "path":          meta.get("path", ""),
            "context":       doc,
            "score":         score,
            "doc_type":      meta.get("doc_type", "General"),
        })

    return results
